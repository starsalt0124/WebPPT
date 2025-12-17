import os
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from utils import px_to_emu, parse_color, parse_rgb_string
from config import PPT_WIDTH_PX, PPT_HEIGHT_PX, TEXT_WIDTH_FACTOR

class PPTRenderer:
    def __init__(self, output_path):
        self.output_path = output_path
        self.prs = Presentation()
        self.prs.slide_width = px_to_emu(PPT_WIDTH_PX)
        self.prs.slide_height = px_to_emu(PPT_HEIGHT_PX)
        self.blank_slide_layout = self.prs.slide_layouts[6]

    def create_text_box(self, slide, el_data):
        """Helper to create a text box from element data"""
        x = px_to_emu(el_data['x'])
        y = px_to_emu(el_data['y'])
        w = px_to_emu(el_data['width'])
        h = px_to_emu(el_data['height'])

        # Store original coordinates for border drawing
        orig_x = x
        orig_w = w

        # Adjust width to prevent wrapping issues
        # And adjust x based on alignment to keep text visually stationary
        w_new = int(w)
        
        if el_data.get('isSingleLine', False):
            w_new = int(w * TEXT_WIDTH_FACTOR)
            
            align = el_data['styles'].get('textAlign', 'left')
            if align == 'center':
                x -= (w_new - w) // 2
            elif align == 'right':
                x -= (w_new - w)
        
        w = w_new

        textbox = slide.shapes.add_textbox(x, y, w, h)
        tf = textbox.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = el_data['text']
        
        # Apply Styles
        font_size_px = float(el_data['styles']['fontSize'].replace('px', ''))
        p.font.size = Pt(font_size_px * 0.75) # Convert px to pt
        
        p.font.color.rgb = parse_rgb_string(el_data['styles']['color'])
        
        # Basic Bold check
        fw = el_data['styles']['fontWeight']
        if (fw.isdigit() and int(fw) >= 600) or fw == 'bold':
            p.font.bold = True
            
        # Font Fallback
        safe_fonts = ["Arial", "Calibri", "Times New Roman", "Microsoft YaHei", "SimHei", "Verdana", "Tahoma"]
        font_family_str = el_data['styles']['fontFamily']
        # Remove quotes and split
        font_families = [f.strip().replace('"', '').replace("'", "") for f in font_family_str.split(',')]
        
        chosen_font = "Arial" # Default
        for f in font_families:
            if f in safe_fonts:
                chosen_font = f
                break
        p.font.name = chosen_font

        # Line Height
        line_height = el_data['styles']['lineHeight']
        if line_height != 'normal':
            if line_height.endswith('px'):
                lh_px = float(line_height.replace('px', ''))
                p.line_spacing = lh_px / font_size_px
            elif line_height.replace('.', '', 1).isdigit():
                 p.line_spacing = float(line_height)

        # Hyperlink
        if el_data.get('href'):
            # Ensure we have a run
            if not p.runs:
                p.add_run()
            r = p.runs[0]
            r.hyperlink.address = el_data['href']

        # Alignment
        align = el_data['styles']['textAlign']
        if align == 'center':
            p.alignment = PP_ALIGN.CENTER
        elif align == 'right':
            p.alignment = PP_ALIGN.RIGHT
        elif align == 'justify':
            p.alignment = PP_ALIGN.JUSTIFY
        else:
            p.alignment = PP_ALIGN.LEFT

        # Vertical Alignment (Text Frame)
        display = el_data['styles']['display']
        align_items = el_data['styles']['alignItems']
        justify_content = el_data['styles']['justifyContent']
        flex_direction = el_data['styles']['flexDirection']

        if display == 'flex' or display == 'inline-flex':
            if flex_direction == 'column':
                if justify_content == 'center':
                    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                elif justify_content == 'flex-end':
                    tf.vertical_anchor = MSO_ANCHOR.BOTTOM
                else:
                    tf.vertical_anchor = MSO_ANCHOR.TOP
            else: # row
                if align_items == 'center':
                    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                elif align_items == 'flex-end':
                    tf.vertical_anchor = MSO_ANCHOR.BOTTOM
                else:
                    tf.vertical_anchor = MSO_ANCHOR.TOP
        else:
            tf.vertical_anchor = MSO_ANCHOR.TOP

        # Background Color (Text Box)
        bg_color = el_data['styles']['backgroundColor']
        if bg_color and bg_color != 'rgba(0, 0, 0, 0)' and bg_color != 'transparent':
            rgb, alpha = parse_color(bg_color)
            textbox.fill.solid()
            textbox.fill.fore_color.rgb = rgb
            if alpha < 1.0:
                textbox.fill.transparency = 1.0 - alpha

        # Border Bottom Handling
        border_bottom_width = el_data['styles'].get('borderBottomWidth')
        if border_bottom_width and border_bottom_width != '0px':
            try:
                width_val = float(border_bottom_width.replace('px', ''))
                if width_val > 0:
                    line = slide.shapes.add_connector(
                        MSO_CONNECTOR.STRAIGHT, 
                        orig_x, y + h, orig_x + orig_w, y + h
                    )
                    line.line.width = Pt(width_val * 0.75)
                    rgb, alpha = parse_color(el_data['styles']['borderBottomColor'])
                    line.line.color.rgb = rgb
            except Exception as e:
                print(f"WARNING: Failed to add bottom border: {e}")

        # Border Left Handling
        border_left_width = el_data['styles'].get('borderLeftWidth')
        if border_left_width and border_left_width != '0px':
            try:
                width_val = float(border_left_width.replace('px', ''))
                if width_val > 0:
                    line = slide.shapes.add_connector(
                        MSO_CONNECTOR.STRAIGHT, 
                        orig_x, y, orig_x, y + h
                    )
                    line.line.width = Pt(width_val * 0.75)
                    rgb, alpha = parse_color(el_data['styles']['borderLeftColor'])
                    line.line.color.rgb = rgb
            except Exception as e:
                print(f"WARNING: Failed to add left border: {e}")

    def create_shape(self, slide, el_data):
        """Helper to create a shape from element data"""
        x = px_to_emu(el_data['x'])
        y = px_to_emu(el_data['y'])
        w = px_to_emu(el_data['width'])
        h = px_to_emu(el_data['height'])
        
        # Determine shape type based on border-radius
        border_radius = el_data['styles']['borderRadius']
        shape_type = MSO_SHAPE.RECTANGLE
        
        if border_radius and border_radius != '0px':
            is_circle = False
            if '50%' in border_radius:
                is_circle = True
            else:
                try:
                    br_val = float(border_radius.replace('px', '').strip())
                    if abs(br_val - el_data['width']/2) < el_data['width']*0.1 and \
                       abs(el_data['width'] - el_data['height']) < el_data['width']*0.1:
                        is_circle = True
                except:
                    pass
            
            if is_circle:
                shape_type = MSO_SHAPE.OVAL
            else:
                shape_type = MSO_SHAPE.ROUNDED_RECTANGLE
        
        shape = slide.shapes.add_shape(shape_type, x, y, w, h)
        
        # Fill Color
        bg_color = el_data['styles']['backgroundColor']
        opacity = float(el_data['styles'].get('opacity', '1'))

        if bg_color and bg_color != 'rgba(0, 0, 0, 0)' and bg_color != 'transparent':
            rgb, alpha = parse_color(bg_color)
            shape.fill.solid()
            shape.fill.fore_color.rgb = rgb
            
            # Combine alpha and opacity
            final_alpha = alpha * opacity
            if final_alpha < 1.0:
                shape.fill.transparency = 1.0 - final_alpha
                # print(f"DEBUG: Shape {el_data.get('id')} transparency set to {1.0 - final_alpha} (alpha={alpha}, opacity={opacity})")
        else:
            shape.fill.background() # No fill
            
        # Border (Line)
        border_width = el_data['styles']['borderTopWidth']
        if border_width and border_width != '0px':
            width_val = float(border_width.replace('px', ''))
            if width_val > 0:
                shape.line.width = Pt(width_val * 0.75)
                rgb, alpha = parse_color(el_data['styles']['borderTopColor'])
                shape.line.color.rgb = rgb
        else:
            shape.line.fill.background() # No line

        # Border Left Handling (for shapes with specific left border)
        border_left_width = el_data['styles'].get('borderLeftWidth')
        if border_left_width and border_left_width != '0px':
            try:
                width_val = float(border_left_width.replace('px', ''))
                if width_val > 0:
                    line = slide.shapes.add_connector(
                        MSO_CONNECTOR.STRAIGHT, 
                        x, y, x, y + h
                    )
                    line.line.width = Pt(width_val * 0.75)
                    rgb, alpha = parse_color(el_data['styles']['borderLeftColor'])
                    line.line.color.rgb = rgb
            except Exception as e:
                print(f"WARNING: Failed to add left border to shape: {e}")

    def create_table(self, slide, el_data):
        """Helper to create a table from element data"""
        x = px_to_emu(el_data['x'])
        y = px_to_emu(el_data['y'])
        w = px_to_emu(el_data['width'])
        h = px_to_emu(el_data['height'])
        
        rows = el_data['rows']
        if not rows:
            return

        num_rows = len(rows)
        num_cols = len(rows[0])
        
        table_shape = slide.shapes.add_table(num_rows, num_cols, x, y, w, h)
        table = table_shape.table
        
        # Set column widths (based on first row)
        for col_idx, cell_data in enumerate(rows[0]):
            table.columns[col_idx].width = px_to_emu(cell_data['width'])

        for row_idx, row_data in enumerate(rows):
            # Set row height (max height of cells in row)
            max_h = max([c['height'] for c in row_data])
            table.rows[row_idx].height = px_to_emu(max_h)
            
            for col_idx, cell_data in enumerate(row_data):
                cell = table.cell(row_idx, col_idx)
                cell.text = cell_data['text']
                
                # Apply Styles to Cell Text
                p = cell.text_frame.paragraphs[0]
                font_size_px = float(cell_data['styles']['fontSize'].replace('px', ''))
                p.font.size = Pt(font_size_px * 0.75)
                p.font.color.rgb = parse_rgb_string(cell_data['styles']['color'])
                
                # Bold check
                fw = cell_data['styles']['fontWeight']
                if (fw.isdigit() and int(fw) >= 600) or fw == 'bold':
                    p.font.bold = True
                    
                # Cell Background
                bg_color = cell_data['styles']['backgroundColor']
                if bg_color and bg_color != 'rgba(0, 0, 0, 0)' and bg_color != 'transparent':
                    rgb, alpha = parse_color(bg_color)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = rgb
                    if alpha < 1.0:
                        cell.fill.transparency = 1.0 - alpha

    def add_slide(self, slide_data, bg_image_path=None):
        slide = self.prs.slides.add_slide(self.blank_slide_layout)
        
        if bg_image_path:
            slide.shapes.add_picture(bg_image_path, 0, 0, self.prs.slide_width, self.prs.slide_height)
        else:
            slide_bg_color = slide_data['backgroundColor']
            background = slide.background
            fill = background.fill
            fill.solid()
            if slide_bg_color == 'rgba(0, 0, 0, 0)' or slide_bg_color == 'transparent':
                 fill.fore_color.rgb = RGBColor(255, 255, 255) # Default to white
            else:
                 fill.fore_color.rgb = parse_rgb_string(slide_bg_color)
        
        return slide

    def add_image_element(self, slide, el_data, image_path, crop_info=None):
        x = px_to_emu(el_data['x'])
        y = px_to_emu(el_data['y'])
        w = px_to_emu(el_data['width'])
        h = px_to_emu(el_data['height'])

        if crop_info:
            # Adjust placement for PPT
            crop_left = crop_info['crop_left']
            crop_top = crop_info['crop_top']
            final_w = crop_info['width']
            final_h = crop_info['height']
            padding = crop_info['padding']
            
            pad_emu = px_to_emu(padding)
            
            # Original intended position (including padding)
            orig_img_x = x - pad_emu
            orig_img_y = y - pad_emu
            
            # New position shifted by crop
            img_x = orig_img_x + px_to_emu(crop_left)
            img_y = orig_img_y + px_to_emu(crop_top)
            
            img_w = px_to_emu(final_w)
            img_h = px_to_emu(final_h)
            
            slide.shapes.add_picture(image_path, img_x, img_y, img_w, img_h)
        else:
            slide.shapes.add_picture(image_path, x, y, w, h)

    def save(self):
        self.prs.save(self.output_path)
